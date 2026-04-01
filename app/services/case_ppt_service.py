from sqlalchemy.orm import Session
from fastapi import HTTPException
from app.models.case_metadata import CaseMetadata
from app.core.config import settings
from app.core.azure_openai import client
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import uuid

class CasePPTService:
    def __init__(self, db: Session):
        self.db = db

    def generate_case_ppt(self, case_id: int) -> str:
        metadata = (
            self.db.query(CaseMetadata)
            .filter(CaseMetadata.case_id == case_id)
            .first()
        )

        if not metadata:
            raise HTTPException(status_code=404, detail="Case metadata not found")

        prompt = f"""
You are a senior litigation lawyer.

Create courtroom-style PPT content.
STRICT FORMAT:

CASE OVERVIEW:
- Bullet points only (max 4)

PARTIES:
- Bullet points (max 3)

BACKGROUND:
- Bullet points (max 4)

LEGAL ISSUES:
- Bullet points (max 4)

EVIDENCE:
- Bullet points (max 4)

IMPORTANT DATES:
- Bullet points (max 4)

STRATEGY & CONCLUSION:
- Bullet points (max 4)

Do NOT use paragraphs.
Do NOT exceed bullet limits.

Case Details:
Court: {metadata.court_name}
Judge: {metadata.judge}
Filing Date: {metadata.filing_date}
Next Court Date: {metadata.next_court_date}
Parties: {metadata.parties}
Attorney: {metadata.attorney}
Strong Evidence: {metadata.strong_evidence}
Description: {metadata.case_description}
"""

        response = client.chat.completions.create(
            model=settings.AZURE_OPENAI_CHAT_DEPLOYMENT,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.25,
            max_tokens=900,
        )

        structured_text = response.choices[0].message.content
        return self._build_styled_ppt(structured_text)

    # ------------------ STYLING CORE ------------------

    def _build_styled_ppt(self, content: str) -> str:
        prs = Presentation()

        sections = self._parse_sections(content)

        for title, bullets in sections.items():
            self._add_lawyer_slide(prs, title, bullets)

        file_path = f"/tmp/case_presentation_{uuid.uuid4().hex}.pptx"
        prs.save(file_path)
        return file_path

    def _parse_sections(self, text: str) -> dict:
        sections = {}
        current_title = None

        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue

            if line.endswith(":"):
                current_title = line.replace(":", "")
                sections[current_title] = []
            elif line.startswith("-") and current_title:
                sections[current_title].append(line[1:].strip())

        return sections

    def _add_lawyer_slide(self, prs, title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # ---------- TITLE ----------
        title_box = slide.shapes.title
        title_box.text = title.upper()

        title_tf = title_box.text_frame
        title_tf.paragraphs[0].font.size = Pt(28)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.name = "Calibri"
        title_tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # ---------- BODY ----------
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for bullet in bullets[:5]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(64, 64, 64)

        tf.paragraphs[0].font.size = Pt(18)
